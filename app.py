import os
import tempfile
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.responses import StreamingResponse
from dotenv import load_dotenv
from PathRAG import PathRAG, QueryParam
from PathRAG.llm import openai_complete, openai_complete_stream
# Additional libraries for file processing
import PyPDF2
import docx2txt
from pptx import Presentation
import openpyxl
from striprtf.striprtf import rtf_to_text
from odf.opendocument import load as load_odf
from odf import teletype
from ebooklib import epub
from bs4 import BeautifulSoup

# Initialize FastAPI with Swagger UI metadata.
app = FastAPI(
    title="PathRAG API",
    description="API for uploading files and querying the RAG system (including streaming responses).",
    version="1.0"
)

# Load environment variables from .env.
load_dotenv()

# Setup a working directory for PathRAG.
WORKING_DIR = os.path.join(os.getcwd(), 'data')
if not os.path.exists(WORKING_DIR):
    os.mkdir(WORKING_DIR)

# Initialize the RAG instance.
rag = PathRAG(
    working_dir=WORKING_DIR,
    llm_model_func=openai_complete,
)

def extract_text_from_file(file: UploadFile) -> str:
    """
    Extract text from an uploaded file.
    Supports many file types including:
    .txt, .md, .pdf, .docx, .pptx, .xlsx, .rtf, .odt, .tex, .epub,
    .html, .htm, .csv, .json, .xml, .yaml, .yml, .log, .conf, .ini,
    .properties, .sql, .bat, .sh, .c, .cpp, .py, .java, .js, .ts,
    .swift, .go, .rb, .php, .css, .scss, .less.
    """
    filename = file.filename
    extension = os.path.splitext(filename)[1].lower()
    file.file.seek(0)

    # Define plain text file extensions.
    plain_text_ext = [
        ".txt", ".md", ".tex", ".csv", ".json", ".xml", ".yaml", ".yml",
        ".log", ".conf", ".ini", ".properties", ".sql", ".bat", ".sh",
        ".c", ".cpp", ".py", ".java", ".js", ".ts", ".swift", ".go",
        ".rb", ".php", ".css", ".scss", ".less"
    ]

    if extension in plain_text_ext:
        try:
            return file.file.read().decode('utf-8', errors='ignore')
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error reading {extension} file: {str(e)}")
    elif extension == ".pdf":
        try:
            file.file.seek(0)
            pdf_reader = PyPDF2.PdfReader(file.file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() or ""
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing PDF file: {str(e)}")
    elif extension == ".docx":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            text = docx2txt.process(tmp_path)
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing DOCX file: {str(e)}")
    elif extension == ".pptx":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            prs = Presentation(tmp_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing PPTX file: {str(e)}")
    elif extension == ".xlsx":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            wb = openpyxl.load_workbook(tmp_path, data_only=True)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    text += row_text + "\n"
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing XLSX file: {str(e)}")
    elif extension == ".rtf":
        try:
            file.file.seek(0)
            content = file.file.read().decode('utf-8', errors='ignore')
            return rtf_to_text(content)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing RTF file: {str(e)}")
    elif extension == ".odt":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".odt") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            doc = load_odf(tmp_path)
            text_content = teletype.extractText(doc)
            os.remove(tmp_path)
            return text_content
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing ODT file: {str(e)}")
    elif extension == ".epub":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".epub") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            book = epub.read_epub(tmp_path)
            text = ""
            for item in book.get_items():
                if item.get_type() == epub.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    text += soup.get_text() + "\n"
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing EPUB file: {str(e)}")
    elif extension in [".html", ".htm"]:
        try:
            file.file.seek(0)
            content = file.file.read().decode('utf-8', errors='ignore')
            soup = BeautifulSoup(content, "html.parser")
            return soup.get_text()
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing HTML file: {str(e)}")
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type.")

@app.post("/upload", summary="Upload a file", description="Upload a file to insert its content into the RAG system.")
async def upload_file(file: UploadFile = File(...)):
    try:
        content = extract_text_from_file(file)
        await rag.ainsert(content)
        return {"message": f"File '{file.filename}' processed and content inserted."}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/query", summary="Query the RAG system", description="Send a query to the RAG system and receive the generated response.")
async def query_rag(query: str):
    try:
        result = await rag.aquery(query, param=QueryParam(mode="hybrid"))
        return {"query": query, "result": result}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/query_stream", summary="Stream Query the RAG system", description="Send a query to the RAG system and stream the generated response.")
async def query_rag_stream(query: str):
    async def stream_generator():
        async for chunk in openai_complete_stream(
            query,
            system_prompt=None,
            history_messages=[],
            keyword_extraction=False,
            param=QueryParam(mode="hybrid")
        ):
            yield chunk
    return StreamingResponse(stream_generator(), media_type="text/plain")
