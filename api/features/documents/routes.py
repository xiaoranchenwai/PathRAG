import logging
from fastapi import Form
from typing import Optional
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File
from sqlalchemy.orm import Session
import os
import shutil
import asyncio
import tempfile
from datetime import datetime
from typing import Optional
from pydantic import BaseModel

from models.database import get_db, Document, User
from api.auth.jwt_handler import get_current_active_user
from .schemas import DocumentResponse, DocumentList, EntitySearchRequest, EntitySearchResponse, SubgraphRequest, SubgraphResponse, EntityNode, RelationshipEdge
from api.features.rag_manager import get_rag_instance, reload_rag_instance

# Additional libraries for file processing
import PyPDF2
import docx2txt
from pptx import Presentation
import openpyxl
from striprtf.striprtf import rtf_to_text
from odf.opendocument import load as load_odf
from odf import teletype
from ebooklib import epub
from bs4 import BeautifulSoup

# Neo4j and GraphML related imports
from neo4j import GraphDatabase
import networkx as nx
import json
import glob

logger = logging.getLogger("PathRAG")

# Create uploads directory if it doesn't exist
UPLOAD_DIR = "./uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

router = APIRouter(
    prefix="/documents",
    tags=["Documents"],
    dependencies=[Depends(get_current_active_user)]
)

# Setup a working directory for PathRAG
WORKING_DIR = os.path.join(os.getcwd(), 'data')
if not os.path.exists(WORKING_DIR):
    os.mkdir(WORKING_DIR)

# Get the RAG instance from the manager
rag = get_rag_instance()

# Neo4j connection schema
class Neo4jConfig(BaseModel):
    uri: str
    username: str
    password: str
    database: Optional[str] = "neo4j"

class Neo4jUploadResponse(BaseModel):
    success: bool
    message: str
    nodes_created: int
    relationships_created: int
    graphml_files_processed: int

def extract_text_from_file(file: UploadFile) -> str:
    """
    Extract text from an uploaded file.
    Supports many file types including:
    .txt, .md, .pdf, .docx, .pptx, .xlsx, .rtf, .odt, .tex, .epub,
    .html, .htm, .csv, .json, .xml, .yaml, .yml, .log, .conf, .ini,
    .properties, .sql, .bat, .sh, .c, .cpp, .py, .java, .js, .ts,
    .swift, .go, .rb, .php, .css, .scss, .less.
    """
    filename = file.filename
    extension = os.path.splitext(filename)[1].lower()
    file.file.seek(0)

    # Define plain text file extensions.
    plain_text_ext = [
        ".txt", ".md", ".tex", ".csv", ".json", ".xml", ".yaml", ".yml",
        ".log", ".conf", ".ini", ".properties", ".sql", ".bat", ".sh",
        ".c", ".cpp", ".py", ".java", ".js", ".ts", ".swift", ".go",
        ".rb", ".php", ".css", ".scss", ".less"
    ]

    if extension in plain_text_ext:
        try:
            return file.file.read().decode('utf-8', errors='ignore')
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error reading {extension} file: {str(e)}")
    elif extension == ".pdf":
        try:
            file.file.seek(0)
            pdf_reader = PyPDF2.PdfReader(file.file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() or ""
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing PDF file: {str(e)}")
    elif extension == ".docx":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            text = docx2txt.process(tmp_path)
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing DOCX file: {str(e)}")
    elif extension == ".pptx":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            prs = Presentation(tmp_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing PPTX file: {str(e)}")
    elif extension == ".xlsx":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            wb = openpyxl.load_workbook(tmp_path, data_only=True)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    text += row_text + "\n"
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing XLSX file: {str(e)}")
    elif extension == ".rtf":
        try:
            file.file.seek(0)
            content = file.file.read().decode('utf-8', errors='ignore')
            return rtf_to_text(content)
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing RTF file: {str(e)}")
    elif extension == ".odt":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".odt") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            doc = load_odf(tmp_path)
            text_content = teletype.extractText(doc)
            os.remove(tmp_path)
            return text_content
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing ODT file: {str(e)}")
    elif extension == ".epub":
        try:
            file.file.seek(0)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".epub") as tmp:
                tmp.write(file.file.read())
                tmp_path = tmp.name
            book = epub.read_epub(tmp_path)
            text = ""
            for item in book.get_items():
                if item.get_type() == epub.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    text += soup.get_text() + "\n"
            os.remove(tmp_path)
            return text
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing EPUB file: {str(e)}")
    elif extension in [".html", ".htm"]:
        try:
            file.file.seek(0)
            content = file.file.read().decode('utf-8', errors='ignore')
            soup = BeautifulSoup(content, "html.parser")
            return soup.get_text()
        except Exception as e:
            raise HTTPException(status_code=400, detail=f"Error processing HTML file: {str(e)}")
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type.")

def load_graphml_files(data_dir: str):
    """
    Load all GraphML files from the data directory.
    Returns a list of NetworkX graph objects.
    """
    graphml_files = glob.glob(os.path.join(data_dir, "*.graphml"))
    graphs = []
    
    for file_path in graphml_files:
        try:
            logger.info(f"Loading GraphML file: {file_path}")
            graph = nx.read_graphml(file_path)
            graphs.append((graph, os.path.basename(file_path)))
        except Exception as e:
            logger.error(f"Error loading GraphML file {file_path}: {str(e)}")
            continue
    
    return graphs

def create_neo4j_session(neo4j_config: Neo4jConfig):
    """
    Create a Neo4j database session.
    """
    try:
        driver = GraphDatabase.driver(
            neo4j_config.uri,
            auth=(neo4j_config.username, neo4j_config.password)
        )
        return driver
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to connect to Neo4j: {str(e)}")

def upload_graph_to_neo4j(driver, graph, filename: str):
    """
    Upload a NetworkX graph to Neo4j database.
    """
    nodes_created = 0
    relationships_created = 0
    
    with driver.session() as session:
        # Clear existing data (optional - you may want to comment this out)
        # session.run("MATCH (n) DETACH DELETE n")
        
        # Create nodes
        for node_id, node_data in graph.nodes(data=True):
            try:
                # Clean node_id by removing quotes if they exist
                clean_node_id = node_id.strip('"')
                
                # Extract properties
                properties = {}
                if 'entity_type' in node_data:
                    properties['entity_type'] = node_data['entity_type'].strip('"')
                if 'description' in node_data:
                    # Handle multiple descriptions separated by <SEP>
                    descriptions = node_data['description'].split('<SEP>')
                    properties['description'] = descriptions[0].strip('"')
                    if len(descriptions) > 1:
                        properties['additional_descriptions'] = [desc.strip('"') for desc in descriptions[1:]]
                if 'source_id' in node_data:
                    properties['source_id'] = node_data['source_id']
                
                # Add metadata
                properties['source_file'] = filename
                properties['node_id'] = clean_node_id
                
                # Create Cypher query
                cypher = """
                MERGE (n:Entity {id: $node_id, source_file: $source_file})
                SET n += $properties
                """
                
                session.run(cypher, {
                    'node_id': clean_node_id,
                    'source_file': filename,
                    'properties': properties
                })
                
                nodes_created += 1
                
            except Exception as e:
                logger.error(f"Error creating node {node_id}: {str(e)}")
                continue
        
        # Create relationships
        for source, target, edge_data in graph.edges(data=True):
            try:
                # Clean node IDs
                clean_source = source.strip('"')
                clean_target = target.strip('"')
                
                # Extract edge properties
                edge_properties = {}
                if 'weight' in edge_data:
                    edge_properties['weight'] = float(edge_data['weight'])
                if 'description' in edge_data:
                    edge_properties['description'] = edge_data['description'].strip('"')
                if 'keywords' in edge_data:
                    edge_properties['keywords'] = edge_data['keywords'].strip('"')
                if 'source_id' in edge_data:
                    edge_properties['source_id'] = edge_data['source_id']
                
                # Add metadata
                edge_properties['source_file'] = filename
                
                # Create relationship
                cypher = """
                MATCH (a:Entity {id: $source, source_file: $source_file})
                MATCH (b:Entity {id: $target, source_file: $source_file})
                MERGE (a)-[r:RELATED]->(b)
                SET r += $properties
                """
                
                session.run(cypher, {
                    'source': clean_source,
                    'target': clean_target,
                    'source_file': filename,
                    'properties': edge_properties
                })
                
                relationships_created += 1
                
            except Exception as e:
                logger.error(f"Error creating relationship {source} -> {target}: {str(e)}")
                continue
    
    return nodes_created, relationships_created

@router.post("/upload", response_model=DocumentResponse, summary="Upload a file", description="Upload a file to insert its content into the RAG system.")
async def upload_file(
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user)
):
    try:
        # Validate file type
        filename = file.filename
        file_extension = os.path.splitext(filename)[1].lower()

        # Map file extensions to content types
        content_type_map = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.md': 'text/markdown',
            '.txt': 'text/plain',
            '.html': 'text/html',
            '.htm': 'text/html',
            '.rtf': 'application/rtf',
            '.odt': 'application/vnd.oasis.opendocument.text',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.epub': 'application/epub+zip'
        }

        content_type = content_type_map.get(file_extension, file.content_type)

        # Create file path
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        file_path = os.path.join(UPLOAD_DIR, f"{current_user.id}_{timestamp}_{filename}")

        # Save file
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Get file size
        file_size = os.path.getsize(file_path)

        # Create document record
        db_document = Document(
            user_id=current_user.id,
            filename=filename,
            content_type=content_type,
            file_path=file_path,
            file_size=file_size,
            status="processing"
        )
        db.add(db_document)
        db.commit()
        db.refresh(db_document)

        # Extract text and process with PathRAG in background
        async def process_document_task():
            # Declare global rag at the beginning of the function
            global rag

            try:
                # Extract text from file
                with open(file_path, "rb") as f:
                    # Create a temporary UploadFile-like object
                    class TempUploadFile:
                        def __init__(self, file_path):
                            self.filename = os.path.basename(file_path)
                            self.file = open(file_path, "rb")

                    temp_file = TempUploadFile(file_path)
                    content = extract_text_from_file(temp_file)
                    temp_file.file.close()

                # Process with PathRAG
                await rag.ainsert(content)

                # Reload the PathRAG instance to make the new document available
                reload_rag_instance()

                # Update the local reference
                rag = get_rag_instance()

                # Update document status
                document = db.query(Document).filter(Document.id == db_document.id).first()
                if document:
                    document.status = "completed"
                    document.processed_at = datetime.now()
                    db.commit()

                logger.info(f"Document processed and PathRAG reloaded successfully.")
            except Exception as e:
                # Update document status on error
                document = db.query(Document).filter(Document.id == db_document.id).first()
                if document:
                    document.status = "failed"
                    document.error_message = str(e)
                    db.commit()

        # Start processing task
        task = asyncio.create_task(process_document_task())

        # Add task to a global set to prevent garbage collection
        if not hasattr(upload_file, 'tasks'):
            upload_file.tasks = set()
        upload_file.tasks.add(task)
        task.add_done_callback(lambda t: upload_file.tasks.remove(t))
        logger.info(f"Document upload completed.")
        return db_document
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/neo4j_upload", response_model=Neo4jUploadResponse, summary="Upload knowledge graph to Neo4j", description="Upload knowledge graph data from /data directory to Neo4j database.")
async def neo4j_upload(
    uri: str ,
    username: str ,
    password: str ,
    database: str = "neo4j",
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_active_user)
):
    # 创建Neo4j配置对象
    neo4j_config = Neo4jConfig(
        uri=uri,
        username=username,
        password=password,
        database=database
    )
    """
    Upload knowledge graph data from GraphML files in /data directory to Neo4j database.
    This endpoint first processes the uploaded document through the regular RAG pipeline,
    then reads all GraphML files from the /data directory and uploads them to Neo4j.
    """
    try:
        # First, process the document through the regular upload pipeline
        logger.info("Starting document processing...")
        
        # Validate file type
        filename = file.filename
        file_extension = os.path.splitext(filename)[1].lower()

        # Map file extensions to content types
        content_type_map = {
            '.pdf': 'application/pdf',
            '.docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            '.md': 'text/markdown',
            '.txt': 'text/plain',
            '.html': 'text/html',
            '.htm': 'text/html',
            '.rtf': 'application/rtf',
            '.odt': 'application/vnd.oasis.opendocument.text',
            '.pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            '.xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            '.epub': 'application/epub+zip'
        }

        content_type = content_type_map.get(file_extension, file.content_type)

        # Create file path
        timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
        file_path = os.path.join(UPLOAD_DIR, f"{current_user.id}_{timestamp}_{filename}")

        # Save file
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)

        # Get file size
        file_size = os.path.getsize(file_path)

        # Create document record
        db_document = Document(
            user_id=current_user.id,
            filename=filename,
            content_type=content_type,
            file_path=file_path,
            file_size=file_size,
            status="processing"
        )
        db.add(db_document)
        db.commit()
        db.refresh(db_document)

        # Process document with PathRAG
        logger.info("Processing document with PathRAG...")
        try:
            # Extract text from file
            with open(file_path, "rb") as f:
                class TempUploadFile:
                    def __init__(self, file_path):
                        self.filename = os.path.basename(file_path)
                        self.file = open(file_path, "rb")

                temp_file = TempUploadFile(file_path)
                content = extract_text_from_file(temp_file)
                temp_file.file.close()
            global rag
            # Process with PathRAG
            await rag.ainsert(content)
            
            # Reload the PathRAG instance
            reload_rag_instance()
            
            rag = get_rag_instance()

            # Update document status
            db_document.status = "completed"
            db_document.processed_at = datetime.now()
            db.commit()

            logger.info("Document processed successfully with PathRAG.")
        except Exception as e:
            logger.error(f"Error processing document with PathRAG: {str(e)}")
            db_document.status = "failed"
            db_document.error_message = str(e)
            db.commit()
            # Continue with Neo4j upload even if RAG processing fails

        # Now process Neo4j upload
        logger.info("Starting Neo4j knowledge graph upload...")
        
        # Test Neo4j connection
        driver = create_neo4j_session(neo4j_config)
        
        # Load GraphML files from /data directory
        graphs = load_graphml_files(WORKING_DIR)
        
        if not graphs:
            raise HTTPException(
                status_code=404, 
                detail=f"No GraphML files found in {WORKING_DIR} directory"
            )
        
        total_nodes_created = 0
        total_relationships_created = 0
        files_processed = 0
        
        # Process each GraphML file
        for graph, graph_filename in graphs:
            logger.info(f"Processing GraphML file: {graph_filename}")
            try:
                nodes_created, relationships_created = upload_graph_to_neo4j(
                    driver, graph, graph_filename
                )
                total_nodes_created += nodes_created
                total_relationships_created += relationships_created
                files_processed += 1
                logger.info(f"Successfully processed {graph_filename}: {nodes_created} nodes, {relationships_created} relationships")
            except Exception as e:
                logger.error(f"Error processing GraphML file {graph_filename}: {str(e)}")
                continue
        
        # Close Neo4j connection
        driver.close()
        
        # Create index on entity id for better performance
        with create_neo4j_session(neo4j_config).session() as session:
            session.run("CREATE INDEX entity_id_index IF NOT EXISTS FOR (n:Entity) ON (n.id)")
            session.run("CREATE INDEX entity_source_file_index IF NOT EXISTS FOR (n:Entity) ON (n.source_file)")
        
        logger.info(f"Neo4j upload completed successfully. Files processed: {files_processed}, Nodes: {total_nodes_created}, Relationships: {total_relationships_created}")
        
        return Neo4jUploadResponse(
            success=True,
            message=f"Successfully uploaded knowledge graph to Neo4j. Document also processed with PathRAG.",
            nodes_created=total_nodes_created,
            relationships_created=total_relationships_created,
            graphml_files_processed=files_processed
        )
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Unexpected error in neo4j_upload: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Internal server error: {str(e)}")

@router.post("/search_entity", response_model=EntitySearchResponse, summary="Search entity by name", description="Search for entity information by name in Neo4j database.")
async def search_entity(
    request: EntitySearchRequest,
    current_user: User = Depends(get_current_active_user)
):
    """
    Search for an entity by name in the Neo4j database and return its information.
    """
    try:
        # Create Neo4j connection
        neo4j_config = Neo4jConfig(
            uri=request.neo4j_uri,
            username=request.neo4j_username,
            password=request.neo4j_password,
            database=request.neo4j_database
        )
        
        driver = create_neo4j_session(neo4j_config)
        
        with driver.session() as session:
            # Search for entity by name (case-insensitive)
            cypher = """
            MATCH (n:Entity)
            WHERE toLower(n.id) = toLower($entity_name)
            RETURN n
            LIMIT 1
            """
            
            result = session.run(cypher, {"entity_name": request.entity_name})
            record = result.single()
            
            if record is None:
                driver.close()
                return EntitySearchResponse(
                    success=False,
                    message=f"Entity '{request.entity_name}' not found in the database.",
                    entity=None
                )
            
            # Extract node properties
            node = record["n"]
            entity_data = {
                "id": node.get("id", ""),
                "entity_type": node.get("entity_type"),
                "description": node.get("description"),
                "additional_descriptions": node.get("additional_descriptions"),
                "source_id": node.get("source_id"),
                "source_file": node.get("source_file"),
                "node_id": node.get("node_id")
            }
            
            entity_node = EntityNode(**entity_data)
            
            driver.close()
            
            return EntitySearchResponse(
                success=True,
                message=f"Entity '{request.entity_name}' found successfully.",
                entity=entity_node
            )
            
    except Exception as e:
        logger.error(f"Error searching for entity '{request.entity_name}': {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Error searching for entity: {str(e)}"
        )

@router.post("/subgraph", response_model=SubgraphResponse, summary="Get subgraph for entity list", description="Return subgraph consisting of specified entities and their relationships.")
async def get_subgraph(
    request: SubgraphRequest,
    current_user: User = Depends(get_current_active_user)
):
    """
    Get a subgraph consisting of the specified entities and all relationships between them.
    """
    try:
        if not request.entity_list:
            raise HTTPException(
                status_code=400,
                detail="Entity list cannot be empty."
            )
        
        # Create Neo4j connection
        neo4j_config = Neo4jConfig(
            uri=request.neo4j_uri,
            username=request.neo4j_username,
            password=request.neo4j_password,
            database=request.neo4j_database
        )
        
        driver = create_neo4j_session(neo4j_config)
        
        with driver.session() as session:
            # First, get all nodes in the entity list
            nodes_cypher = """
            MATCH (n:Entity)
            WHERE toLower(n.id) IN [entity IN $entity_list | toLower(entity)]
            RETURN n
            """
            
            nodes_result = session.run(nodes_cypher, {"entity_list": request.entity_list})
            nodes = []
            
            for record in nodes_result:
                node = record["n"]
                entity_data = {
                    "id": node.get("id", ""),
                    "entity_type": node.get("entity_type"),
                    "description": node.get("description"),
                    "additional_descriptions": node.get("additional_descriptions"),
                    "source_id": node.get("source_id"),
                    "source_file": node.get("source_file"),
                    "node_id": node.get("node_id")
                }
                nodes.append(EntityNode(**entity_data))
            
            # Get relationships between the entities in the list
            relationships_cypher = """
            MATCH (n1:Entity)-[r:RELATED]-(n2:Entity)
            WHERE toLower(n1.id) IN [entity IN $entity_list | toLower(entity)]
              AND toLower(n2.id) IN [entity IN $entity_list | toLower(entity)]
            RETURN n1.id as source, n2.id as target, r
            """
            
            relationships_result = session.run(relationships_cypher, {"entity_list": request.entity_list})
            relationships = []
            
            for record in relationships_result:
                rel = record["r"]
                relationship_data = {
                    "source": record["source"],
                    "target": record["target"],
                    "weight": rel.get("weight"),
                    "description": rel.get("description"),
                    "keywords": rel.get("keywords"),
                    "source_id": rel.get("source_id"),
                    "source_file": rel.get("source_file")
                }
                relationships.append(RelationshipEdge(**relationship_data))
            
            driver.close()
            
            return SubgraphResponse(
                success=True,
                message=f"Subgraph retrieved successfully with {len(nodes)} nodes and {len(relationships)} relationships.",
                nodes=nodes,
                relationships=relationships,
                nodes_count=len(nodes),
                relationships_count=len(relationships)
            )
            
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error retrieving subgraph: {str(e)}")
        raise HTTPException(
            status_code=500,
            detail=f"Error retrieving subgraph: {str(e)}"
        )

@router.get("/", response_model=DocumentList)
async def get_documents(db: Session = Depends(get_db), current_user: User = Depends(get_current_active_user)):
    documents = db.query(Document).filter(Document.user_id == current_user.id).all()
    return {"documents": documents}

@router.get("/{document_id}", response_model=DocumentResponse)
async def get_document(document_id: int, db: Session = Depends(get_db), current_user: User = Depends(get_current_active_user)):
    document = db.query(Document).filter(Document.id == document_id, Document.user_id == current_user.id).first()
    if document is None:
        raise HTTPException(status_code=404, detail="Document not found")
    return document

@router.get("/{document_id}/status", response_model=DocumentResponse)
async def get_document_status(document_id: int, db: Session = Depends(get_db), current_user: User = Depends(get_current_active_user)):
    # Check if document exists and belongs to user
    document = db.query(Document).filter(Document.id == document_id, Document.user_id == current_user.id).first()
    if document is None:
        raise HTTPException(status_code=404, detail="Document not found")

    # Return the document with its current status
    return document

@router.post("/reload", response_model=dict)
async def reload_documents(current_user: User = Depends(get_current_active_user)):
    """
    Reload the PathRAG instance to recognize newly uploaded documents.
    This is useful after uploading documents to make them available for querying
    without restarting the server.
    """
    # Declare global rag at the beginning of the function
    global rag

    try:
        # Reload the PathRAG instance
        reload_rag_instance()

        # Update the local reference
        rag = get_rag_instance()

        return {
            "success": True,
            "message": "PathRAG instance reloaded successfully. New documents are now available for querying."
        }
    except Exception as e:
        raise HTTPException(
            status_code=500,
            detail=f"Failed to reload PathRAG instance: {str(e)}"
        )