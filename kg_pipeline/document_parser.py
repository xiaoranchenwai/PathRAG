import os
import tempfile

import PyPDF2
import docx2txt
import openpyxl
from bs4 import BeautifulSoup
from ebooklib import epub
from odf import teletype
from odf.opendocument import load as load_odf
from pptx import Presentation
from striprtf.striprtf import rtf_to_text


PLAIN_TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".tex",
    ".csv",
    ".json",
    ".xml",
    ".yaml",
    ".yml",
    ".log",
    ".conf",
    ".ini",
    ".properties",
    ".sql",
    ".bat",
    ".sh",
    ".c",
    ".cpp",
    ".py",
    ".java",
    ".js",
    ".ts",
    ".swift",
    ".go",
    ".rb",
    ".php",
    ".css",
    ".scss",
    ".less",
}


def extract_text_from_path(file_path: str) -> str:
    filename = os.path.basename(file_path)
    extension = os.path.splitext(filename)[1].lower()

    if extension in PLAIN_TEXT_EXTENSIONS:
        with open(file_path, "rb") as file_handle:
            return file_handle.read().decode("utf-8", errors="ignore")

    if extension == ".pdf":
        with open(file_path, "rb") as file_handle:
            pdf_reader = PyPDF2.PdfReader(file_handle)
            return "".join(page.extract_text() or "" for page in pdf_reader.pages)

    if extension == ".docx":
        return docx2txt.process(file_path)

    if extension == ".pptx":
        presentation = Presentation(file_path)
        text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    if extension == ".xlsx":
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        text = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) for cell in row if cell is not None)
                if row_text:
                    text.append(row_text)
        return "\n".join(text)

    if extension == ".rtf":
        with open(file_path, "rb") as file_handle:
            content = file_handle.read().decode("utf-8", errors="ignore")
        return rtf_to_text(content)

    if extension == ".odt":
        with tempfile.NamedTemporaryFile(delete=False, suffix=".odt") as tmp:
            with open(file_path, "rb") as file_handle:
                tmp.write(file_handle.read())
            tmp_path = tmp.name
        document = load_odf(tmp_path)
        os.remove(tmp_path)
        return teletype.extractText(document)

    if extension == ".epub":
        book = epub.read_epub(file_path)
        text = []
        for item in book.get_items():
            if item.get_type() == epub.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                text.append(soup.get_text())
        return "\n".join(text)

    if extension in {".html", ".htm"}:
        with open(file_path, "rb") as file_handle:
            content = file_handle.read().decode("utf-8", errors="ignore")
        soup = BeautifulSoup(content, "html.parser")
        return soup.get_text()

    raise ValueError(f"Unsupported file type: {extension}")
